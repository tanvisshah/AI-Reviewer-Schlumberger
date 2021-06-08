from flask_login import login_required
from flask import Flask, render_template, request,session,flash,url_for,redirect,send_from_directory


# this is needed after this for migration to work properly
from app.models import *
# from models import User,Uploads
from app import app
from app import login_manager
from app.forms import *

import string
from nltk.stem import WordNetLemmatizer
import nltk
nltk.download('averaged_perceptron_tagger')
import os
import sys
from werkzeug.utils import secure_filename
import logging
from logging import Formatter, FileHandler
from pdf2image import convert_from_path
import subprocess
import glob
from pptx import Presentation
from pptx.dml.color import RGBColor
from gingerit.gingerit import GingerIt
import numpy as np
from pptx.enum.text import MSO_UNDERLINE
import language_tool_python
# from flask_weasyprint import HTML, render_pdf

# determine if application is a script file or frozen exe
if getattr(sys, 'frozen', False):
    application_path = os.path.dirname(sys.executable)
    print("application_path",application_path)
    print("Meipass",sys._MEIPASS)
    application_path = sys._MEIPASS
elif __file__:
    application_path = os.path.dirname(__file__)
    print("application_path",application_path)


@login_manager.user_loader
def load_user(id):
    return User.query.get(int(id))


@app.shell_context_processor
def make_shell_context():
    return {'db': db, 'User': User}

# Automatically tear down SQLAlchemy.
'''
@app.teardown_request
def shutdown_session(exception=None):
    db_session.remove()
'''

# Login required decorator.
'''
def login_required(test):
    @wraps(test)
    def wrap(*args, **kwargs):
        if 'logged_in' in session:
            return test(*args, **kwargs)
        else:
            flash('You need to login first.')
            return redirect(url_for('login'))
    return wrap
'''
#----------------------------------------------------------------------------#
# Controllers.
#----------------------------------------------------------------------------#

from sqlalchemy import desc


@app.route('/upload_presentation', methods=['GET', 'POST'])
def upload_daily_file():
    if request.method == 'POST':
        file = request.files['file']

        projectName = request.form['projectnamedropdown']

        if projectName == "Other":
            projectName = request.form['projectname']

        print(projectName)

        session['projectName'] = projectName
        # print("FILES:",files[0].filename)
        if file.filename != '':
            session['current_file_name'] = file.filename
            user = User.query.filter_by(username=current_user.username).first()
            
            print("User Id",user.id)

            projectQuery = Uploads.query.filter_by(id=user.id,project_name=projectName).first()

            # print("Does path exist?",os.path.isdir("data/"+current_user.username))
            userFolder = os.path.join(application_path,"data/"+current_user.username)
            if not os.path.isdir(userFolder):
                os.makedirs(userFolder)
            
            if projectQuery == None:
                projectFolder = userFolder+"/"+"".join(projectName.split())
                os.makedirs(projectFolder)
                u = Uploads(id=user.id,project_name=projectName,version_number=1,uploaded_filename=file.filename)
                db.session.add(u)
                db.session.commit()
                projectFolder = projectFolder+"/1"
                os.makedirs(projectFolder)

            else:
                projectQuery = Uploads.query.filter_by(id=user.id,project_name=projectName).order_by(desc(Uploads.version_number)).limit(1).first()
                newVersionNumber = int(projectQuery.version_number) + 1
                
                # Make new version folder
                projectFolder = userFolder+"/"+"".join(projectName.split())
                projectFolder = projectFolder+"/"+str(newVersionNumber)
                os.makedirs(projectFolder)

                u = Uploads(id=user.id,project_name=projectName,version_number=newVersionNumber,uploaded_filename=file.filename)
                db.session.add(u)
                db.session.commit()
            print(projectQuery)

            file.save(os.path.join(projectFolder,secure_filename("".join(file.filename.split()))))
                # current_file_name = f

        print("current isjbeiwudcv: ",session['current_file_name'])
        return redirect("/process_file") # Redirect to home


@app.route('/index')
def index():

    return redirect("/")



@app.route('/')
@login_required
def home():

    user = User.query.filter_by(username=current_user.username).first()
    projectQuery = Uploads.query.filter_by(id=user.id).group_by(Uploads.project_name).all()

    project_names = [project.project_name for project in projectQuery]
    project_names.append("Other")
    print(project_names)
    print("HERE",session)
    
    return render_template('pages/placeholder.home.html',current_projects = project_names)


def get_modified_slides(current_file_name):
    tool = language_tool_python.LanguageTool('en-US')
    
    user = User.query.filter_by(username=current_user.username).first()
    projectQuery = Uploads.query.filter_by(id=user.id,project_name=session['projectName']).order_by(desc(Uploads.version_number)).limit(1).first()
    session['current_version'] = str(projectQuery.version_number)
    folder_path = os.path.join(application_path,"data/"+current_user.username+"/"+"".join(projectQuery.project_name.split())+"/"+str(projectQuery.version_number))

    print("Presentation Uploaded Filepath",folder_path)
    prs = Presentation(folder_path+"/"+"".join(projectQuery.uploaded_filename.split()))
    
    error_types=[]
    num_errors = np.zeros(len(prs.slides))
    error_slides=[]
    

    # for i,slide in enumerate(prs.slides):
    #     print("Errors found in slide: ",i)
    #     corrections=[]
    #     for shape in slide.shapes:
    #         if shape.has_text_frame:
    #             frame = shape.text_frame
    #             for p in frame.paragraphs:
    #                 text = p.text
    #                 matches = tool.check(text)
                    
    #                 if len(matches)!=0:
    #                     for err in matches:
    #                         num_errors[i] += 1
    #                         # print(err)
    #                         corrections.append(err)
    #                         text_to_correct = text[err.offset:err.offset+err.errorLength]
    #                         error_types.append(err.ruleIssueType)
    #                         if p.text.find(text_to_correct)!=-1:
    #                             p.font.underline = True
    #                             p.font.color.rgb = RGBColor(0xff, 0x00, 0x00)
    #                             p.font.underline = MSO_UNDERLINE.WAVY_LINE
    #                         print()
    #     corrections.append(i)
    #     error_slides.append(corrections)                   
    #     print()


    for i,slide in enumerate(prs.slides):
        print("* Errors found in slide: "+str(i)+" *")
        print()
        corrections=[]
        for shape in slide.shapes:
            if shape.has_text_frame:
                frame = shape.text_frame
                for p in frame.paragraphs:
                    text = p.text
                    
                    t = text.split()
                    for w in range(len(t)):
                        t[w] = t[w].translate(str.maketrans('', '', string.punctuation))
                    # print("Error Here:",t)


                    tokens_tag = nltk.pos_tag([i for i in t if i])
                    matches = tool.check(text)

                    if len(matches)!=0:
                        for err in matches:
                            flag = False
                            text_to_correct = text[err.offset:err.offset+err.errorLength]
                            #if not present in tokens tag, check it
                            if not(any(text_to_correct in k for k in tokens_tag)):
                                flag = True
                            #if present in tokens tag, check for NNP. If not NNP, check it
                            else:
                                for j in tokens_tag:
                                    if j[0]==text_to_correct:
                                        if j[1] != 'NNP' and j[1]!='NNPS':
                                            flag = True
                            

                            if flag == True:
                                num_errors[i]+=1
                                corrections.append("Error found in: "+text+"\n"+
                                                   "Erroneous word: "+text_to_correct+"\n"+
                                                    "Message: "+err.message+"\n"+
                                                    "Suggestions: "+str(err.replacements[:10]))
                                print("Error found in: "+text)
                                print("Erroneous word: "+text_to_correct)
                                print("Message: "+err.message)
                                print("Suggestions: "+str(err.replacements[:10]))
                                error_types.append(err.ruleIssueType)
                                if p.text.find(text_to_correct)!=-1:
                                    p.font.underline = True
                                    p.font.color.rgb = RGBColor(0xff, 0x00, 0x00)
                                    p.font.underline = MSO_UNDERLINE.WAVY_LINE
                                print()
        corrections.append(i)
        error_slides.append(corrections)

    prs.save(folder_path+"/modified_"+projectQuery.uploaded_filename)

    return error_slides,error_types,num_errors,folder_path,projectQuery.uploaded_filename



@app.route('/process_file')
@login_required
def process_file():
    
    # Store Pdf with convert_from_path function
    # print("Current File Name",current_file_name)

    # current_file_name = glob.glob("uploaded_files/*.pptx")[0].split("/")[-1]

    print("Current: ",session['current_file_name'])

    error_slides,error_types,num_errors,path,filename = get_modified_slides(session['current_file_name'])

    modified_filename = "modified_"+filename

    print("Error Slides:",error_slides,type(error_slides))
    # print("Path:",os.getcwd()+"/uploaded_files/modified_"+session['current_file_name'])
    from sys import platform
    if platform == "linux" or platform == "linux2":
        subprocess.run(["soffice", "--headless" ,"--convert-to", "pdf", path+"/"+modified_filename,"--outdir",path])
    elif platform == "darwin":
        subprocess.run(["/Applications/LibreOffice.app/Contents/MacOS/soffice", "--headless" ,"--convert-to", "pdf", path+"/"+modified_filename,"--outdir",path])
    
    # session['current_file_name'].split('.')[:-1]

    images = convert_from_path(path+"/"+modified_filename.split('.')[0]+'.pdf')
    
    slides_with_errors = []
    error_text = []

    for slide in error_slides:
        if len(slide)>1:
            slides_with_errors.append(slide[-1])
            error_text.append(slide[:-1])

    # files = glob.glob('static/converted_files/*')
    # for f in files:
    #     os.remove(f)


    for i in slides_with_errors:
        if(i<10):
            images[i].save(path+"/"+ str('0')+str(i) +'.png', 'png')
        else:
            images[i].save(path+"/"+ str(i) +'.png', 'png')


    # for i in range(len(images)):
        # Save pages as images in the pdf
        # images[i].save('static/converted_files/page'+ str(i) +'.jpg', 'JPEG')

    # files = glob.glob("static/converted_files/*")
    # print(files)

    # files = os.listdir('static/converted_files')
    # files.sort()
    # files = ['converted_files/' +file for file in files]

    # print(files)
    # return render_template('pages/placeholder.postprocess.html',files=files)

    summary_list=[]

    (unique, counts) = np.unique(error_types, return_counts=True)
    frequencies = np.asarray((unique, counts)).T
    # print("Total number of errors found: "+str(len(error_types)))

    summary_list.append(["Total Errors:",str(len(error_types))])
    print()

    for i in range(len(frequencies)):
        # print("Total number of "+frequencies[i][0]+" errors: "+frequencies[i][1])
        summary_list.append([frequencies[i][0]+" errors: ",frequencies[i][1]])

    print()
    errors_per_slide = []
    for i in range(len(num_errors)):
        if num_errors[i]!=0:
            # print("Number of errors in slide "+str(i+1)+": "+str(num_errors[i]))
            errors_per_slide.append(["Errors in slide "+str(i+1)+": ",str(num_errors[i])])

    np.save(path+"/error_text"+filename,error_text)
    np.save(path+"/slides_with_errors"+filename,slides_with_errors)
    np.save(path+"/summary_list"+filename,summary_list)
    np.save(path+"/errors_per_slide"+filename,errors_per_slide)

    session['error_text'] = error_text
    session['slides_with_errors'] = slides_with_errors
    session['summary_list'] = summary_list
    

    return redirect("/show_corrections/{}/{}".format( "".join(session['projectName'].split()),session['current_version'] ))

# @app.route('/get_pdf/<projectname>/<id>')
# @login_required
# def get_pdf(projectname,id):
#     # print("HERE",session)
#     # return redirect("/")
#     # return render_pdf(url_for('show_corrections'))

#     # print("Sessifffffon",session)

#     path = "data/"+current_user.username+"/"+"".join(projectname.split())+"/"+id
#     files = glob.glob(path+"/*.png")
#     # print("Using Glob",files)
#     files.sort()

#     # files = os.listdir('static/converted_files')
#     # files.sort()
#     files = [file for file in files]


#     # nltk.download('wordnet')
#     pres_path = glob.glob(path+"/*.pptx")[0]
#     print("presentation path",pres_path)
#     # prs = Presentation('uploaded_files/'+"".join(session['current_file_name'].split()))
#     prs = Presentation(pres_path)
#     lemmatizer = WordNetLemmatizer()

#     num_slides = len(prs.slides)
#     title_words = []
    
#     print("Slides:",prs.slides[0].shapes.title.text)
#     for i in range(num_slides):
#         try:
#             t = prs.slides[i].shapes.title.text
#             title_words += (t.strip().split())
#         except:
#             pass
#     title_words = [lemmatizer.lemmatize(x.lower()) for x in title_words]
#     required = ['introduction', 'conclusion', 'summary', 'methodology', 'assumption', 'reference']
#     missing_titles = []
#     for req in required:
#         if req not in title_words:
#             missing_titles.append((req,0))
#         else:
#             missing_titles.append((req,1))

#     # print("Current Files to Display:",files)
#     # print("Error text:",session['error_text'])

#     error_text = np.load(glob.glob(path+"/error_text*.npy")[0],allow_pickle=True)
#     slides_with_errors = np.load(glob.glob(path+"/slides_with_errors*.npy")[0],allow_pickle=True)
#     summary_list = np.load(glob.glob(path+"/summary_list*.npy")[0],allow_pickle=True)
#     errors_per_slide = np.load(glob.glob(path+"/errors_per_slide*.npy")[0],allow_pickle=True)
    
#     print("ID_VERSION",id)
#     html = render_template('pages/placeholder.postprocess.html', data = zip(files,error_text,slides_with_errors,errors_per_slide),missing_titles=missing_titles,summary_list=summary_list,projectname=projectname,version_id=id)

#     # html = render_template('pages/placeholder.postprocess.html', data = zip(files,session['error_text'],session['slides_with_errors']),missing_titles=missing_titles,summary_list=session['summary_list'])



#     return render_pdf(HTML(string=html))

@app.route('/projects')
@login_required
def projects():

    user = User.query.filter_by(username=current_user.username).first()

    print("User Id",user.id)

    projectQuery = Uploads.query.filter_by(id=user.id).group_by(Uploads.project_name).all()

    project_names = [project.project_name for project in projectQuery]
    versions=[]
    
    for name in project_names:
        temp=[]
        projectQuery = Uploads.query.filter_by(id=user.id,project_name=name).all()
        print("printing this:",projectQuery)
        for result in projectQuery:
            print("Results",result.version_number)
            temp.append(result.version_number)
        versions.append(temp)

    print(versions)

    return render_template('pages/projects.html',project_names=list(zip(project_names,versions)))


@app.route('/slide_images/<path:path>')
def slide_images(path):
    print("PATHHERE",path)
    import urllib.parse
    path = urllib.parse.unquote(path)
    # print("Printing parsedpath",path)
    # print("Application Path:",application_path)
    if application_path == "./app":
        return send_from_directory('/app/',path, conditional=True)
    else:
        return send_from_directory('/',path, conditional=True)



@app.route('/show_corrections/<projectname>/<id>')
@login_required
def show_corrections(projectname,id):
    # print("Sessifffffon",session)

    path = os.path.join(application_path,"data/"+current_user.username+"/"+"".join(projectname.split())+"/"+id)
    files = glob.glob(path+"/*.png")
    # print("Using Glob",files)
    files.sort()

    # files = os.listdir('static/converted_files')
    # files.sort()
    files = [file for file in files]


    # nltk.download('wordnet')
    pres_path = glob.glob(path+"/*.pptx")[0]
    print("presentation path",pres_path)
    # prs = Presentation('uploaded_files/'+"".join(session['current_file_name'].split()))
    prs = Presentation(pres_path)
    lemmatizer = WordNetLemmatizer()

    num_slides = len(prs.slides)
    title_words = []
    
    print("Slides:",prs.slides[0].shapes.title.text)
    for i in range(num_slides):
        try:
            t = prs.slides[i].shapes.title.text
            title_words += (t.strip().split())
        except:
            pass
    title_words = [lemmatizer.lemmatize(x.lower()) for x in title_words]
    required = ['introduction', 'conclusion', 'summary', 'methodology', 'assumption', 'reference']
    missing_titles = []
    for req in required:
        if req not in title_words:
            missing_titles.append((req,0))
        else:
            missing_titles.append((req,1))

    


    # print("Current Files to Display:",files)
    # print("Error text:",session['error_text'])

    error_text = np.load(glob.glob(path+"/error_text*.npy")[0],allow_pickle=True)
    slides_with_errors = np.load(glob.glob(path+"/slides_with_errors*.npy")[0],allow_pickle=True)
    summary_list = np.load(glob.glob(path+"/summary_list*.npy")[0],allow_pickle=True)
    errors_per_slide = np.load(glob.glob(path+"/errors_per_slide*.npy")[0],allow_pickle=True)

    # print("SUMMARY LIST:",summary_list)

    

    return render_template('pages/placeholder.postprocess.html', data = zip(files,error_text,slides_with_errors,errors_per_slide),missing_titles=missing_titles,summary_list=summary_list,projectname=projectname,version_id=id)

    # return render_template('pages/placeholder.postprocess.html',files=files)


@app.route('/about')
@login_required
def about():
    return render_template('pages/placeholder.about.html')


from flask_login import current_user, login_user, logout_user


@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('index'))



@app.route('/login',methods=["POST","GET"])
def login():

    if current_user.is_authenticated:
        print("Current Logged in User:",current_user.username)
        
        return redirect(url_for('index'))

    form = LoginForm()
    print("HERE")
    if form.validate_on_submit():
        flash("Login requested for user {}".format(form.username.data))

        user = User.query.filter_by(username=form.username.data).first()
        if user is None or not user.check_password(form.password.data):
            flash('Invalid username or password')
            return redirect(url_for('login'))

        login_user(user, remember=form.remember_me.data)

        session['name'] = form.username.data
        print(session['name'])
        # print("HERE")
        return redirect("/")
    # flash('ERRORS')
    print("Couldn't alidate",form.errors)
    return render_template('forms/login.html', form=form)

# app.route('/login',methods=["POST","GET"])
# def login():
#     form = LoginForm()
#     if form.validate_on_submit():
#         flash("Login requested for user {}".format(form.username.data))
#         session['name'] = form.username.data
#         print(session['name'])
#         print("HERE")
#         return redirect("/")
        
#     return render_template("forms/login.html",title="Sign In",form=form)



@app.route('/register',methods=["GET","POST"])
def register():
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    form = RegisterForm()
    # flash("FLASHING MESSAGES")
    if form.validate_on_submit():
        # due to circular import 
        # https://stackoverflow.com/questions/26126989/python-flask-sqlalchemy-cannot-import-from-models

        # print("Submit successful")
        user = User(username=form.username.data, email=form.email.data)

        user.set_password(form.password.data)
        db.session.add(user)
        db.session.commit()
        flash('Congratulations, you are now a registered user!')
        return redirect(url_for('login'))
    return render_template('forms/register.html', form=form)

# app.route('/register',methods=["GET","POST"])
# def register():
#     if current_user.is_authenticated:
#         return redirect(url_for('index'))
#     form = RegisterForm()
#     print("Pirnt this")
#     if form.validate_on_submit():
#         # due to circular import 
#         # https://stackoverflow.com/questions/26126989/python-flask-sqlalchemy-cannot-import-from-models

#         print("Submit successful")
#         user = User(username=form.username.data, email=form.email.data)

#         user.set_password(form.password.data)
#         db.session.add(user)
#         db.session.commit()
#         flash('Congratulations, you are now a registered user!')
#         return redirect(url_for('login'))
#     return render_template('forms/register.html', form=form)
#     # return render_template(url_for("login.html"),form=form)




@app.route('/forgot')
def forgot():
    form = ForgotForm(request.form)
    return render_template('forms/forgot.html', form=form)

# Error handlers.


@app.errorhandler(500)
def internal_error(error):
    #db_session.rollback()
    return render_template('errors/500.html'), 500


@app.errorhandler(404)
def not_found_error(error):
    return render_template('errors/404.html'), 404

if not app.debug:
    file_handler = FileHandler('error.log')
    file_handler.setFormatter(
        Formatter('%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]')
    )
    app.logger.setLevel(logging.INFO)
    file_handler.setLevel(logging.INFO)
    app.logger.addHandler(file_handler)
    app.logger.info('errors')